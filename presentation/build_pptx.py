"""Build Fischer AI-Lösungen presentation (German, easy language)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ── Brand colours ─────────────────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0x1E, 0x3A, 0x8A)   # deep brand blue
BLUE_MID   = RGBColor(0x3B, 0x82, 0xF6)   # accent blue
ORANGE     = RGBColor(0xF9, 0x73, 0x16)   # accent orange
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_TEXT  = RGBColor(0x47, 0x55, 0x69)

# Slide dimensions (widescreen 16:9)
W = Cm(33.87)
H = Cm(19.05)

# ── Slide data (pairs of ideas per slide) ─────────────────────────────────────
SLIDES: list[dict] = [
    {
        "kind": "title",
        "title": "Mehr Kunden.\nWeniger Arbeit.",
        "subtitle": "Wie schlauer Computereinsatz das Geschäft von\nFischer Entrümpelungen noch besser macht.",
    },
    {
        "kind": "pair",
        "left": {
            "icon": "💬",
            "title": "Rund-um-die-Uhr\nKundenberater",
            "bullets": [
                "Beantwortet Fragen auf Deutsch, Englisch, Türkisch und Arabisch",
                "Gibt sofort einen Preishinweis – auch abends und am Wochenende",
                "Bucht Termine und leitet Anfragen ans Team weiter",
                "Kein Anruf geht mehr verloren",
            ],
        },
        "right": {
            "icon": "📷",
            "title": "Angebot per\nFoto",
            "bullets": [
                "Kunde schickt einfach Fotos von der Wohnung",
                "Das Programm schätzt den Aufwand automatisch",
                "Erstes Angebot in wenigen Minuten – nicht erst nach Stunden",
                "Mehr Aufträge, weniger Telefonzeit",
            ],
        },
    },
    {
        "kind": "pair",
        "left": {
            "icon": "📋",
            "title": "Schlauere\nKundenverwaltung",
            "bullets": [
                "Anfragen per E-Mail, WhatsApp und Telefon kommen automatisch ins System",
                "Dringende Jobs werden sofort nach oben gesetzt",
                "Das Team sieht auf einen Blick, was zu tun ist",
                "Weniger Papierkram, mehr Zeit für Kunden",
            ],
        },
        "right": {
            "icon": "⭐",
            "title": "Mehr gute\nBewertungen",
            "bullets": [
                "Nach jedem Auftrag: automatische freundliche Nachricht an den Kunden",
                "Bittet um eine Google-Bewertung zum richtigen Zeitpunkt",
                "Antwortet professionell auf Bewertungen – kein Aufwand fürs Team",
                "Bessere Noten auf Google = mehr Anrufe",
            ],
        },
    },
    {
        "kind": "pair",
        "left": {
            "icon": "🗺️",
            "title": "Clevere\nRoutenplanung",
            "bullets": [
                "Alle Jobs des Tages werden auf dem kürzesten Weg eingeplant",
                "Weniger Fahrzeit und Sprit – mehr Aufträge pro Tag",
                "Das Team weiß morgens genau, was wann wo ansteht",
                "Auch bei kurzfristigen Änderungen passt sich der Plan an",
            ],
        },
        "right": {
            "icon": "📝",
            "title": "Werbetexte\nvon selbst",
            "bullets": [
                "Erstellt automatisch neue Texte für die Website und Google",
                "Schreibt Facebook-Posts und Saisonangebote auf Knopfdruck",
                "Vorher-Nachher-Geschichten aus echten Aufträgen – ganz einfach",
                "Immer sichtbar online – ohne extra Marketingabteilung",
            ],
        },
    },
    {
        "kind": "pair",
        "left": {
            "icon": "📸",
            "title": "Schönere Fotos\nautomatisch",
            "bullets": [
                "Auftragsfotos werden automatisch verbessert und aufgehellt",
                "Persönliche Daten auf Bildern werden entfernt",
                "Kurze Videos für Social Media – ohne Videoagentur",
                "Professionelles Auftreten macht Vertrauen",
            ],
        },
        "right": {
            "icon": "📞",
            "title": "Telefon-\nAssistent",
            "bullets": [
                "Nimmt jeden Anruf entgegen – auch wenn alle beschäftigt sind",
                "Notiert Name, Adresse und was der Kunde braucht",
                "Bucht Termine direkt im Kalender",
                "Nur komplizierte Fälle kommen ans Team",
            ],
        },
    },
    {
        "kind": "pair",
        "left": {
            "icon": "📊",
            "title": "Alles auf\neinen Blick",
            "bullets": [
                "Eine Übersicht für Aufträge, Umsatz und Kundenanfragen",
                "Zeigt, welche Stadtteile am meisten buchen",
                "Gibt Tipps, wo Werbung sich am meisten lohnt",
                "Entscheidungen einfach und schnell treffen",
            ],
        },
        "right": {
            "icon": "📅",
            "title": "Schlauer\nPersonalplan",
            "bullets": [
                "Erkennt automatisch, wann viel los ist (Frühling, Umzugssaison)",
                "Plant Personal und Fahrzeuge rechtzeitig ein",
                "Keine Überraschungen mehr in der Hochsaison",
                "Mehr Aufträge annehmen – ohne Chaos",
            ],
        },
    },
    {
        "kind": "closing",
        "title": "Der nächste Schritt",
        "points": [
            "Kurzes Gespräch: Was passt am besten zu Fischer?",
            "Kleiner Test mit einem echten Werkzeug – kostenlos",
            "Schritt für Schritt einführen – ohne alles auf einmal zu ändern",
        ],
        "footer": "Wir machen das gemeinsam. Einfach. Sicher. Erfolgreich.",
    },
]


# ── Helpers ────────────────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(
    slide, left, top, width, height,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = WHITE,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_rect(slide, left, top, width, height, color: RGBColor, transparency: int = 0) -> None:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_bullet_box(
    slide, left, top, width, height,
    bullets: list[str],
    font_size: int = 14,
    color: RGBColor = GRAY_TEXT,
) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"✔  {bullet}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


# ── Slide builders ────────────────────────────────────────────────────────────

def _build_title_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, BLUE_DARK)

    # Decorative orange bar left
    _add_rect(slide, Cm(0), Cm(0), Cm(1.2), H, ORANGE)

    # White block centre
    _add_rect(slide, Cm(3), Cm(3.5), Cm(26), Cm(11.5), WHITE)

    # Company name tag
    _add_text_box(slide, Cm(3.5), Cm(1.2), Cm(20), Cm(1.5),
                  "Fischer Entrümpelungen", 13, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

    # Main title
    _add_text_box(slide, Cm(3.5), Cm(3.8), Cm(25), Cm(5),
                  data["title"], 38, bold=True, color=BLUE_DARK, align=PP_ALIGN.LEFT)

    # Subtitle
    _add_text_box(slide, Cm(3.5), Cm(10), Cm(24), Cm(4),
                  data["subtitle"], 17, bold=False, color=GRAY_TEXT, align=PP_ALIGN.LEFT)

    # Bottom tagline
    _add_text_box(slide, Cm(3), Cm(16), Cm(28), Cm(1.5),
                  "KI-Lösungen · einfach erklärt", 13, bold=False, color=BLUE_MID, align=PP_ALIGN.CENTER)


def _build_pair_slide(prs: Presentation, data: dict, slide_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, GRAY_LIGHT)

    # Top header bar
    _add_rect(slide, Cm(0), Cm(0), W, Cm(2.2), BLUE_DARK)
    _add_text_box(slide, Cm(1), Cm(0.3), Cm(28), Cm(1.8),
                  "Was der Computer für Fischer tun kann", 15, bold=True,
                  color=WHITE, align=PP_ALIGN.LEFT)
    _add_text_box(slide, Cm(29), Cm(0.3), Cm(4), Cm(1.8),
                  f"{slide_num}", 15, bold=False, color=BLUE_MID, align=PP_ALIGN.RIGHT)

    # ── LEFT card ────────────────────────────────────────────────────────
    lx, ly = Cm(0.8), Cm(2.5)
    lw, lh = Cm(15.5), Cm(15.8)
    _add_rect(slide, lx, ly, lw, lh, WHITE)
    # Orange top accent
    _add_rect(slide, lx, ly, lw, Cm(0.35), ORANGE)

    left = data["left"]
    _add_text_box(slide, lx + Cm(0.5), ly + Cm(0.5), Cm(3), Cm(2),
                  left["icon"], 36, color=BLUE_DARK, align=PP_ALIGN.LEFT)
    _add_text_box(slide, lx + Cm(0.5), ly + Cm(2.5), lw - Cm(1), Cm(2.5),
                  left["title"], 22, bold=True, color=BLUE_DARK, align=PP_ALIGN.LEFT)
    _add_bullet_box(slide, lx + Cm(0.5), ly + Cm(5.2), lw - Cm(1), Cm(10),
                    left["bullets"], font_size=14, color=GRAY_TEXT)

    # ── RIGHT card ───────────────────────────────────────────────────────
    rx, ry = Cm(17.2), Cm(2.5)
    rw, rh = Cm(15.9), Cm(15.8)
    _add_rect(slide, rx, ry, rw, rh, WHITE)
    _add_rect(slide, rx, ry, rw, Cm(0.35), BLUE_MID)

    right = data["right"]
    _add_text_box(slide, rx + Cm(0.5), ry + Cm(0.5), Cm(3), Cm(2),
                  right["icon"], 36, color=BLUE_DARK, align=PP_ALIGN.LEFT)
    _add_text_box(slide, rx + Cm(0.5), ry + Cm(2.5), rw - Cm(1), Cm(2.5),
                  right["title"], 22, bold=True, color=BLUE_DARK, align=PP_ALIGN.LEFT)
    _add_bullet_box(slide, rx + Cm(0.5), ry + Cm(5.2), rw - Cm(1), Cm(10),
                    right["bullets"], font_size=14, color=GRAY_TEXT)


def _build_closing_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, BLUE_DARK)

    _add_rect(slide, Cm(0), Cm(0), Cm(1.2), H, ORANGE)

    _add_text_box(slide, Cm(3), Cm(1.5), Cm(28), Cm(2),
                  data["title"], 34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    _add_rect(slide, Cm(3), Cm(4), Cm(27), Cm(0.08), BLUE_MID)

    for i, point in enumerate(data["points"]):
        _add_text_box(slide, Cm(3), Cm(4.8 + i * 2.8), Cm(27), Cm(2.4),
                      f"{i+1}.  {point}", 19, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    # Closing badge
    _add_rect(slide, Cm(3), Cm(15.2), Cm(27), Cm(2.5), ORANGE)
    _add_text_box(slide, Cm(3.3), Cm(15.4), Cm(26.4), Cm(2.1),
                  data["footer"], 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def build() -> Path:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    pair_num = 0
    for slide_data in SLIDES:
        kind = slide_data["kind"]
        if kind == "title":
            _build_title_slide(prs, slide_data)
        elif kind == "pair":
            pair_num += 1
            _build_pair_slide(prs, slide_data, pair_num)
        elif kind == "closing":
            _build_closing_slide(prs, slide_data)

    out = Path(__file__).parent / "Fischer_KI_Loesungen.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    return out


if __name__ == "__main__":
    build()
